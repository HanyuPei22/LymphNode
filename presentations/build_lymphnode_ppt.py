from __future__ import annotations

import math
import textwrap
from pathlib import Path

import fitz
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path("/work/h0pei001/workspace/LymphNode_DSN")
PAPER = ROOT / "paper" / "DSN_2026_Hanyu_LymphNode.pdf"
FIG_DIR = ROOT / "paper" / "figures"
OUT_DIR = ROOT / "presentations"
ASSET_DIR = OUT_DIR / "assets" / "rendered_figures"
PPTX_OUT = OUT_DIR / "LymphNode_DSN_17min_presentation.pptx"
NOTES_OUT = OUT_DIR / "LymphNode_DSN_17min_speaker_notes.md"


COLORS = {
    "ink": RGBColor(28, 35, 44),
    "muted": RGBColor(89, 101, 116),
    "line": RGBColor(211, 218, 226),
    "panel": RGBColor(246, 248, 251),
    "blue": RGBColor(39, 94, 176),
    "green": RGBColor(29, 145, 99),
    "orange": RGBColor(225, 126, 58),
    "red": RGBColor(184, 69, 69),
    "purple": RGBColor(110, 88, 176),
    "white": RGBColor(255, 255, 255),
    "dark": RGBColor(17, 24, 39),
}


SLIDES = [
    {
        "title": "LymphNode",
        "subtitle": "A plug-and-play access control method for deep neural networks",
        "kicker": "DSN 2026 | 17-minute presentation draft",
        "body": [
            "Goal: protect deployed DNNs before theft, not only prove ownership after theft.",
            "Core idea: default-deny inference with feature-space credentials and sparse adversarial neutralization.",
        ],
        "time": "0:00-0:40",
        "notes": "Open with the concrete problem: valuable DNNs must often be deployed outside a controlled cloud. LymphNode is an access-control layer inside the model, not a watermark used after a theft has already happened.",
    },
    {
        "title": "Talk Structure",
        "layout": "timeline",
        "items": [
            ("1", "Problem", "Why edge deployment creates oracle access"),
            ("2", "Prior Work", "Watermarking, locking, output poisoning"),
            ("3", "Design", "Feature-domain credential + GSUAP"),
            ("4", "Evidence", "Neutralization, overhead, attacks, robustness"),
        ],
        "time": "0:40-1:10",
        "notes": "Set expectations. About one third of the talk is motivation and related work, one third is design, and one third is results.",
    },
    {
        "title": "Why This Problem Matters",
        "layout": "two_column",
        "left_title": "Deployment pressure",
        "left": [
            "Models are high-value IP.",
            "Healthcare, finance, and edge deployments often cannot rely on remote APIs.",
            "On-device or on-premise inference improves latency and data locality.",
        ],
        "right_title": "Security consequence",
        "right": [
            "Adversaries can query the model at high volume.",
            "Black-box outputs can leak decision behavior.",
            "A copied model can avoid licensing and enable downstream attacks.",
        ],
        "visual": "edge_loop",
        "time": "1:10-1:55",
        "notes": "Frame the conflict: deployment moves models closer to users and attackers. That is good for latency and privacy, but it weakens cloud-style access control.",
    },
    {
        "title": "Threat: Oracle Access Enables Model Theft",
        "layout": "attack_loop",
        "bullets": [
            "Model extraction: train a substitute from query-response pairs.",
            "Model inversion: use outputs to reconstruct sensitive training information.",
            "Edge deployment removes practical rate limits and server-side query filtering.",
        ],
        "sources": [
            "Tramer et al., USENIX Security 2016",
            "Orekondy et al., CVPR 2019",
            "Fredrikson et al., CCS 2015",
        ],
        "time": "1:55-2:45",
        "notes": "Use this slide to explain why the attacker does not need weights. The attacker only needs repeated queries and output probabilities or labels.",
    },
    {
        "title": "Related Work: Three Defense Families",
        "layout": "three_cards",
        "cards": [
            ("Passive provenance", "Watermarks / fingerprints verify ownership after theft.", "Gap: does not prevent unauthorized execution."),
            ("Model locking", "Cryptographic or structural authorization gates model utility.", "Gap: retraining, data access, or high runtime cost."),
            ("Output poisoning", "Perturb predictions to make stolen surrogates worse.", "Gap: assumes cloud-side interposition layer."),
        ],
        "time": "2:45-3:35",
        "notes": "This is the related-work taxonomy. The important transition is that LymphNode borrows the access-control objective from active defenses, but tries to remove their setup and deployment friction.",
    },
    {
        "title": "Where Existing Active Defenses Struggle",
        "layout": "matrix",
        "columns": ["Method type", "Original data", "Retrain", "Post-hoc", "Edge fit"],
        "rows": [
            ["Deep-Lock / crypto", "No", "No", "Yes", "High latency"],
            ["SSAT / ActiveGuard", "Often yes", "Yes", "No", "Setup heavy"],
            ["AdvParams", "No", "No", "Yes", "Parameter keying"],
            ["PP / CIP", "No", "No", "API layer", "Cloud-oriented"],
            ["LymphNode", "No or tiny", "No", "Yes", "Constant cost"],
        ],
        "time": "3:35-4:20",
        "notes": "Do not over-claim that prior work is useless. The claim is narrower: none of them simultaneously gives post-hoc deployment, low setup cost, low overhead, and edge compatibility.",
    },
    {
        "title": "Contribution in One Sentence",
        "layout": "statement",
        "statement": "LymphNode turns a trained DNN into a default-deny model: unauthorized inputs hit a sparse feature-space perturbation, while authorized inputs carry a hidden credential that cancels the perturbation.",
        "chips": ["Post-hoc", "Data-efficient", "Feature-space", "Active defense", "Low overhead"],
        "time": "4:20-5:00",
        "notes": "This is the anchor sentence. Repeat the two flows: no credential means neutralized output; valid credential means antidote and clean utility.",
    },
    {
        "title": "Threat Model and Deployment Roles",
        "layout": "roles",
        "roles": [
            ("Model owner", "Trains model, builds LymphNode plugin, issues keys"),
            ("Edge operator", "Deploys protected model in trusted runtime"),
            ("Authorized user", "Submits input with feature-domain credential"),
            ("Adversary", "Has architecture knowledge and oracle queries, but no key or feature access"),
        ],
        "time": "5:00-5:45",
        "notes": "Clarify the assumption: LymphNode protects against oracle-access abuse under runtime integrity. It does not claim to solve arbitrary write access to model files.",
    },
    {
        "title": "System Overview",
        "layout": "image",
        "image": "framework",
        "caption": "The plugin sits after the early feature layer: sparse GSUAP is active by default; a verified credential triggers inverse noise.",
        "time": "5:45-6:40",
        "notes": "Walk left to right. Clean image plus credential becomes authorized input. The checkpoint verifies feature bits. If matched, inverse noise cancels GSUAP. Otherwise, the target DNN sees corrupted features.",
    },
    {
        "title": "Feature-Domain Credential",
        "layout": "credential",
        "bullets": [
            "Secret key: 32-bit binary pattern embedded in selected early-layer features.",
            "Verification extracts quantized feature bits, not a visible patch.",
            "Collision probability is approximately 2^-32 under balanced bits.",
            "Authorized input is created by a small pixel perturbation that satisfies the feature pattern.",
        ],
        "time": "6:40-7:30",
        "notes": "Emphasize why this is not a BadNets-style visible trigger. The credential lives in a feature representation and uses fine-grained quantized bits.",
    },
    {
        "title": "GSUAP: Neutralize by Touching Critical Channels",
        "layout": "pipeline",
        "steps": [
            ("1", "Rank channels", "Use weight-gradient sensitivity on a small calibration set."),
            ("2", "Build sparse mask", "Select the top-k decision-critical channels."),
            ("3", "Optimize universal delta", "Train a static feature-space perturbation under an L-infinity bound."),
            ("4", "Deploy default-on", "Normal users receive delta; VIP users receive delta plus antidote."),
        ],
        "time": "7:30-8:30",
        "notes": "The method has two phases: decide where to intervene, then learn one universal perturbation. During inference, the perturbation is just a tensor addition.",
    },
    {
        "title": "Inference Logic",
        "layout": "flow",
        "left": ["Input x", "Conv1 / patch embedding", "Checkpoint extracts 32 bits"],
        "branches": [
            ("No match", "Add masked GSUAP", "Utility neutralized"),
            ("Match", "Add GSUAP + antidote", "High-fidelity output"),
        ],
        "time": "8:30-9:10",
        "notes": "This slide is useful if the audience missed the figure. Make the decision logic explicit: same backbone, different feature-space treatment based on credential verification.",
    },
    {
        "title": "Artifact and Code Mapping",
        "layout": "code_map",
        "items": [
            ("Feature verification", "src/models/control_model.py"),
            ("UAP/GSUAP injection", "src/models/uap_control_model.py"),
            ("GD-UAP optimization", "src/training/gd_uap_trainer.py"),
            ("Channel selection", "src/selection/weight_gradient_based.py"),
            ("Main reproduction", "reproduce.py --exp1"),
        ],
        "time": "9:10-9:45",
        "notes": "For an artifact-aware DSN audience, this slide helps show that the method is implemented as a clear pipeline, not just a conceptual figure.",
    },
    {
        "title": "Experimental Setup",
        "layout": "setup",
        "bullets": [
            "Datasets: CIFAR-10, MNIST, SVHN; additional cross-dataset and CelebA attack studies.",
            "Models: ResNet-18/50, ViT-Tiny/Small, plus DenseNet and AlexNet in the paper table.",
            "Baselines: Gaussian noise, SUAP, PP, CIP, BadNets, Blended, Passport.",
            "Metrics: unauthorized accuracy down, VIP accuracy up, overhead, attack success, robustness.",
        ],
        "time": "9:45-10:25",
        "notes": "State the evaluation contract before showing results. Lower unauthorized accuracy is better; high VIP accuracy means normal service for authorized users.",
    },
    {
        "title": "Main Result: Lock Out Unauthorized Queries",
        "layout": "bars",
        "bars": [
            ("Gaussian", 85.4, COLORS["orange"]),
            ("SUAP", 72.0, COLORS["purple"]),
            ("GSUAP", 13.6, COLORS["green"]),
            ("VIP", 94.5, COLORS["blue"]),
        ],
        "caption": "ResNet-18 / CIFAR-10 / 60% channel ratio. Unauthorized accuracy drops near random guessing while VIP accuracy stays high.",
        "time": "10:25-11:25",
        "notes": "This is the most important quantitative slide. The contrast is simple: naive noise barely locks the model, SUAP helps, GSUAP makes the model almost useless to unauthorized users while keeping VIP accuracy at 94.5%.",
    },
    {
        "title": "Efficiency Frontier",
        "layout": "image",
        "image": "efficiency_comparison",
        "caption": "GSUAP gives larger accuracy gap per modified channel, especially on robust architectures.",
        "time": "11:25-12:05",
        "notes": "Explain the efficiency metric as protection benefit per channel touched. The point is not only that GSUAP works, but that it works efficiently.",
    },
    {
        "title": "Runtime Cost is Constant and Small",
        "layout": "overhead",
        "metrics": [
            ("ResNet-18", "1.70 -> 2.74 ms", "+1.0 ms latency", "-14.8% throughput"),
            ("ViT-Tiny", "1.54 -> 2.60 ms", "+1.1 ms latency", "-6.7% throughput"),
            ("Storage/FLOPs", "No meaningful increase", "Params unchanged", "Element-wise operations"),
        ],
        "time": "12:05-12:45",
        "notes": "This is the practical deployment argument. LymphNode does not run per-query optimization; it adds a fixed sparse tensor operation.",
    },
    {
        "title": "Protection Against Real Attack Pipelines",
        "layout": "attack_results",
        "left_title": "Knockoff model extraction",
        "left": [
            "No defense: 85.24% surrogate accuracy at 50K queries.",
            "PP/CIP: 45.66% / 48.25%.",
            "LymphNode: 15.28%, close to random guessing.",
        ],
        "right_title": "GMI model inversion",
        "right": [
            "Clean target: Acc-1 82.71%, Acc-5 91.32%.",
            "LymphNode: Acc-1 4.17%, Acc-5 11.22%.",
            "Latency: 0.088 ms vs PP 0.308 ms.",
        ],
        "time": "12:45-13:50",
        "notes": "This connects the method back to the opening threat model: if labels and probabilities are corrupted in feature space, the data collected by extraction and inversion attacks becomes uninformative.",
    },
    {
        "title": "Low Data Requirement",
        "layout": "image",
        "image": "exp3_visualization",
        "caption": "GSUAP reaches robust lockout with roughly 50-100 calibration samples at practical channel ratios.",
        "time": "13:50-14:30",
        "notes": "This is a key differentiator. Prior active defenses often need original data or retraining. LymphNode can initialize with a very small calibration set.",
    },
    {
        "title": "Cross-Dataset Adaptivity",
        "layout": "mini_table",
        "columns": ["Source -> Target", "Ratio", "Accuracy Gap", "Fooling Rate"],
        "rows": [
            ["CIFAR-10 -> CIFAR-100", "0.4", "69.98%", "64.13%"],
            ["CIFAR-10 -> CIFAR-100", "0.6", "70.66%", "75.02%"],
            ["CIFAR-10 -> STL-10", "0.6", "46.06%", "50.20%"],
            ["CIFAR-10 -> STL-10", "0.8", "51.31%", "55.11%"],
        ],
        "time": "14:30-15:05",
        "notes": "If the original data is unavailable, a public surrogate can still give useful protection when visual statistics are related. Acknowledge STL-10 is harder due to domain and resolution shift.",
    },
    {
        "title": "Robustness and Stealth",
        "layout": "four_cards",
        "cards": [
            ("Imperceptible key", "PSNR 56.64 dB, SSIM 0.9990, LPIPS 0.0011"),
            ("Forgery resistance", "0.0% success for linear residual and U-Net attacks"),
            ("Fine-tuning resistance", "Recovered clean accuracy only 52.46% after 50 epochs"),
            ("JPEG recovery", "Q=60 reaches 85.6% ASR after iterative embedding"),
        ],
        "time": "15:05-15:55",
        "notes": "Summarize robustness without drowning in details. The credential is hard to see, hard to forge, and the lock is not easily removed by normal fine-tuning.",
    },
    {
        "title": "Limitations and Honest Scope",
        "layout": "two_column",
        "left_title": "Assumptions",
        "left": [
            "Runtime integrity is preserved.",
            "Attacker lacks the secret key and intermediate feature access.",
            "The model owner can distribute credentials out of band.",
        ],
        "right_title": "Limitations",
        "right": [
            "Software-only plugin needs hardening against full file tampering.",
            "Aggressive INT8 quantization may distort fine-grained credentials.",
            "Future work: obfuscation, hardware-backed integrity, quantization-aware embedding.",
        ],
        "time": "15:55-16:30",
        "notes": "This slide makes the defense credible. It is an access-control layer under a clear deployment model, not a universal protection against arbitrary physical compromise.",
    },
    {
        "title": "Takeaway",
        "layout": "takeaway",
        "headline": "LymphNode shifts DNN IP protection from proof-after-theft to prevention-before-use.",
        "points": [
            "Default-deny feature-space control.",
            "Post-hoc integration without retraining.",
            "Strong unauthorized neutralization with high VIP fidelity.",
            "Low setup cost and roughly 1 ms inference overhead.",
        ],
        "time": "16:30-17:00",
        "notes": "End with the sentence the audience should remember. If time remains, point to the artifact and the figures for reproducibility.",
    },
    {
        "title": "Backup: Useful Sources and Artifact Map",
        "layout": "sources",
        "sources": [
            ("Model extraction", "Tramer et al. 2016; Knockoff Nets 2019"),
            ("Watermarking / provenance", "Uchida 2017; Adi 2018; DeepSigns 2019; SoK 2022"),
            ("Active authorization", "Deep-Lock 2020; AdvParams 2021; ActiveGuard/SSAT"),
            ("UAP foundation", "Universal adversarial perturbations; GD-UAP"),
            ("Artifact", "README.md, reproduce.py, src/models, src/training, src/selection"),
        ],
        "time": "backup",
        "notes": "Keep this as a backup slide for Q&A or if someone asks where the implementation maps to the paper.",
    },
]


def trim_white(path: Path) -> None:
    img = Image.open(path).convert("RGB")
    bg = Image.new("RGB", img.size, (255, 255, 255))
    diff = ImageChops.difference(img, bg)
    bbox = diff.getbbox()
    if bbox:
        pad = 20
        left = max(0, bbox[0] - pad)
        top = max(0, bbox[1] - pad)
        right = min(img.size[0], bbox[2] + pad)
        bottom = min(img.size[1], bbox[3] + pad)
        img.crop((left, top, right, bottom)).save(path)


def render_pdf_figures() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    rendered = {}
    for pdf in FIG_DIR.glob("*.pdf"):
        out = ASSET_DIR / f"{pdf.stem}.png"
        doc = fitz.open(pdf)
        page = doc[0]
        pix = page.get_pixmap(matrix=fitz.Matrix(3.0, 3.0), alpha=False)
        pix.save(out)
        doc.close()
        trim_white(out)
        rendered[pdf.stem] = out
    rendered["overview"] = ROOT / "docs" / "overview.png"
    return rendered


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def set_line(shape, color, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(slide, text, x, y, w, h, size=22, color=None, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_bullets(slide, bullets, x, y, w, h, size=20, color=None, gap=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.word_wrap = True
    for i, item in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(gap)
    return box


def add_header(slide, title, idx):
    add_text(slide, title, 0.62, 0.28, 10.7, 0.55, size=25, bold=True)
    add_text(slide, f"{idx:02d}", 12.15, 0.32, 0.45, 0.35, size=12, color=COLORS["muted"], align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.92), Inches(12.05), Inches(0.02))
    set_fill(line, COLORS["line"])


def add_footer(slide, text="LymphNode | DSN 2026"):
    add_text(slide, text, 0.62, 7.08, 5.6, 0.22, size=8, color=COLORS["muted"])


def add_panel(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    set_fill(shape, fill or COLORS["panel"])
    set_line(shape, line or COLORS["line"], 1)
    return shape


def add_chip(slide, text, x, y, fill):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.45), Inches(0.33))
    shape.adjustments[0] = 0.22
    set_fill(shape, fill)
    shape.line.color.rgb = fill
    add_text(slide, text, x + 0.05, y + 0.055, 1.35, 0.18, size=9, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)


def add_picture_fit(slide, img_path: Path, x, y, w, h):
    img = Image.open(img_path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    return slide.shapes.add_picture(str(img_path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def draw_arrow(slide, x1, y1, x2, y2, color=None, width=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color or COLORS["blue"]
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def slide_title(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 250, 252)
    add_text(slide, s["kicker"], 0.75, 0.65, 5.8, 0.35, size=12, color=COLORS["green"], bold=True)
    add_text(slide, s["title"], 0.72, 1.23, 7.6, 0.8, size=43, bold=True)
    add_text(slide, s["subtitle"], 0.76, 2.08, 7.8, 0.5, size=19, color=COLORS["muted"])
    add_bullets(slide, s["body"], 0.85, 3.1, 5.6, 1.2, size=18, gap=4)
    add_panel(slide, 7.65, 1.05, 4.75, 4.65, fill=COLORS["white"])
    add_text(slide, "Default-deny DNN", 8.0, 1.35, 4.0, 0.35, size=19, bold=True)
    add_text(slide, "Unauthorized", 8.1, 2.2, 1.6, 0.3, size=12, color=COLORS["red"], bold=True)
    add_text(slide, "neutralized utility", 9.9, 2.2, 1.7, 0.3, size=12, color=COLORS["red"])
    draw_arrow(slide, 8.1, 2.6, 11.1, 2.6, COLORS["red"], 2.5)
    add_text(slide, "Authorized", 8.1, 3.4, 1.6, 0.3, size=12, color=COLORS["green"], bold=True)
    add_text(slide, "clean fidelity", 9.9, 3.4, 1.7, 0.3, size=12, color=COLORS["green"])
    draw_arrow(slide, 8.1, 3.8, 11.1, 3.8, COLORS["green"], 2.5)
    add_text(slide, "Feature-space checkpoint", 8.0, 4.55, 3.6, 0.35, size=16, bold=True)
    add_text(slide, "GSUAP active by default; credential triggers antidote.", 8.0, 4.95, 3.8, 0.5, size=12, color=COLORS["muted"])
    add_footer(slide)


def slide_timeline(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    x0, y = 0.85, 2.05
    for i, (num, title, desc) in enumerate(s["items"]):
        x = x0 + i * 3.0
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.72), Inches(0.72))
        set_fill(circle, [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]][i])
        add_text(slide, num, x, y + 0.14, 0.72, 0.25, size=14, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
        if i < len(s["items"]) - 1:
            draw_arrow(slide, x + 0.82, y + 0.36, x + 2.55, y + 0.36, COLORS["line"], 2)
        add_text(slide, title, x - 0.1, y + 1.0, 2.2, 0.3, size=18, bold=True)
        add_text(slide, desc, x - 0.1, y + 1.42, 2.25, 0.7, size=13, color=COLORS["muted"])
    add_text(slide, "Suggested timing: 5 min background, 4 min design, 7 min evidence, 1 min close.", 1.0, 5.9, 10.8, 0.45, size=18, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_two_column(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_panel(slide, 0.75, 1.38, 5.55, 4.95, fill=COLORS["white"])
    add_panel(slide, 6.95, 1.38, 5.55, 4.95, fill=COLORS["white"])
    add_text(slide, s["left_title"], 1.1, 1.72, 4.8, 0.35, size=20, bold=True, color=COLORS["blue"])
    add_bullets(slide, s["left"], 1.12, 2.25, 4.7, 2.2, size=17, gap=4)
    add_text(slide, s["right_title"], 7.3, 1.72, 4.8, 0.35, size=20, bold=True, color=COLORS["red"])
    add_bullets(slide, s["right"], 7.32, 2.25, 4.7, 2.2, size=17, gap=4)
    if s.get("visual") == "edge_loop":
        add_text(slide, "cloud API", 1.25, 5.05, 1.1, 0.25, size=11, color=COLORS["muted"])
        add_text(slide, "edge model", 3.05, 5.05, 1.25, 0.25, size=11, color=COLORS["muted"])
        add_text(slide, "attacker", 4.75, 5.05, 1.0, 0.25, size=11, color=COLORS["muted"])
        draw_arrow(slide, 2.25, 5.18, 3.0, 5.18, COLORS["line"], 2)
        draw_arrow(slide, 4.35, 5.18, 4.75, 5.18, COLORS["line"], 2)
    add_footer(slide)


def slide_attack_loop(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    nodes = [
        ("Adversary", 0.95, 2.0, COLORS["red"]),
        ("Queries", 3.55, 1.15, COLORS["orange"]),
        ("Victim DNN", 6.0, 2.0, COLORS["blue"]),
        ("Outputs", 3.55, 3.95, COLORS["purple"]),
        ("Stolen / inverted model", 8.95, 2.0, COLORS["dark"]),
    ]
    for text, x, y, color in nodes:
        add_panel(slide, x, y, 2.05, 0.9, fill=COLORS["white"], line=color)
        add_text(slide, text, x + 0.1, y + 0.27, 1.85, 0.25, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    draw_arrow(slide, 2.95, 2.28, 3.55, 1.62, COLORS["orange"], 2.5)
    draw_arrow(slide, 5.62, 1.62, 6.0, 2.28, COLORS["orange"], 2.5)
    draw_arrow(slide, 6.0, 2.72, 5.62, 4.4, COLORS["purple"], 2.5)
    draw_arrow(slide, 3.55, 4.4, 2.95, 2.72, COLORS["purple"], 2.5)
    draw_arrow(slide, 8.05, 2.45, 8.95, 2.45, COLORS["red"], 2.5)
    add_bullets(slide, s["bullets"], 1.0, 5.3, 7.2, 1.1, size=15, color=COLORS["ink"], gap=2)
    add_text(slide, "Sources: " + "; ".join(s["sources"]), 1.0, 6.62, 10.8, 0.25, size=8, color=COLORS["muted"])
    add_footer(slide)


def slide_three_cards(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    colors = [COLORS["blue"], COLORS["orange"], COLORS["purple"]]
    for i, (name, desc, gap) in enumerate(s["cards"]):
        x = 0.8 + i * 4.05
        add_panel(slide, x, 1.55, 3.5, 4.7, fill=COLORS["white"], line=colors[i])
        add_text(slide, name, x + 0.25, 1.9, 2.95, 0.35, size=18, bold=True, color=colors[i])
        add_text(slide, desc, x + 0.25, 2.55, 2.95, 0.75, size=14, color=COLORS["ink"])
        add_text(slide, gap, x + 0.25, 4.35, 2.95, 0.8, size=15, color=COLORS["red"], bold=True)
    add_footer(slide)


def slide_matrix(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    x, y = 0.75, 1.45
    col_w = [2.9, 1.7, 1.45, 1.45, 3.0]
    row_h = 0.62
    for c, header in enumerate(s["columns"]):
        add_panel(slide, x + sum(col_w[:c]), y, col_w[c], row_h, fill=COLORS["dark"], line=COLORS["dark"])
        add_text(slide, header, x + sum(col_w[:c]) + 0.06, y + 0.18, col_w[c] - 0.1, 0.16, size=9, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    for r, row in enumerate(s["rows"]):
        ry = y + row_h * (r + 1)
        fill = RGBColor(241, 247, 244) if "LymphNode" in row[0] else COLORS["white"]
        for c, val in enumerate(row):
            add_panel(slide, x + sum(col_w[:c]), ry, col_w[c], row_h, fill=fill, line=COLORS["line"])
            color = COLORS["green"] if "LymphNode" in row[0] else COLORS["ink"]
            add_text(slide, val, x + sum(col_w[:c]) + 0.06, ry + 0.18, col_w[c] - 0.1, 0.16, size=9, color=color, bold=("LymphNode" in row[0]), align=PP_ALIGN.CENTER)
    add_text(slide, "Positioning: LymphNode aims for active prevention without original-data dependence, retraining, or per-query server logic.", 1.0, 5.85, 10.8, 0.5, size=17, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_statement(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_panel(slide, 1.05, 1.7, 11.1, 2.35, fill=COLORS["white"], line=COLORS["green"])
    add_text(slide, s["statement"], 1.45, 2.05, 10.25, 1.25, size=27, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    for i, chip in enumerate(s["chips"]):
        add_chip(slide, chip, 1.3 + i * 2.25, 4.75, [COLORS["blue"], COLORS["green"], COLORS["orange"], COLORS["purple"], COLORS["dark"]][i])
    add_footer(slide)


def slide_roles(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    positions = [(0.85, 1.55), (7.0, 1.55), (0.85, 4.15), (7.0, 4.15)]
    colors = [COLORS["blue"], COLORS["green"], COLORS["purple"], COLORS["red"]]
    for i, (role, desc) in enumerate(s["roles"]):
        x, y = positions[i]
        add_panel(slide, x, y, 5.45, 1.55, fill=COLORS["white"], line=colors[i])
        add_text(slide, role, x + 0.28, y + 0.28, 4.7, 0.3, size=18, bold=True, color=colors[i])
        add_text(slide, desc, x + 0.28, y + 0.78, 4.7, 0.45, size=14, color=COLORS["muted"])
    draw_arrow(slide, 3.5, 3.15, 6.5, 3.15, COLORS["line"], 2)
    add_text(slide, "Secure out-of-band key distribution; oracle-only attacker.", 4.0, 3.27, 5.3, 0.25, size=12, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_image(prs, s, idx, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_picture_fit(slide, images[s["image"]], 0.85, 1.25, 11.65, 4.9)
    add_text(slide, s["caption"], 0.9, 6.35, 11.3, 0.42, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_credential(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_bullets(slide, s["bullets"], 0.85, 1.35, 5.5, 4.8, size=17, gap=4)
    add_panel(slide, 7.0, 1.45, 4.8, 4.25, fill=COLORS["white"], line=COLORS["blue"])
    add_text(slide, "32-bit feature key", 7.35, 1.78, 4.0, 0.35, size=19, bold=True, color=COLORS["blue"])
    bit_y = 2.55
    for i in range(32):
        row = i // 8
        col = i % 8
        bit = "1" if i in {0, 3, 5, 9, 12, 17, 18, 21, 25, 30} else "0"
        color = COLORS["green"] if bit == "1" else COLORS["line"]
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.35 + col * 0.42), Inches(bit_y + row * 0.42), Inches(0.31), Inches(0.31))
        set_fill(rect, color)
        rect.line.color.rgb = COLORS["white"]
        add_text(slide, bit, 7.35 + col * 0.42, bit_y + row * 0.06 + row * 0.42, 0.31, 0.12, size=8, color=COLORS["dark"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Extracted from selected channels and spatial locations after early feature extraction.", 7.35, 4.65, 3.75, 0.5, size=12, color=COLORS["muted"])
    add_footer(slide)


def slide_pipeline(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    for i, (num, title, desc) in enumerate(s["steps"]):
        x = 0.9 + i * 3.05
        add_panel(slide, x, 2.0, 2.45, 2.7, fill=COLORS["white"], line=[COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]][i])
        add_chip(slide, num, x + 0.35, 2.32, [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]][i])
        add_text(slide, title, x + 0.28, 3.0, 1.9, 0.35, size=16, bold=True)
        add_text(slide, desc, x + 0.28, 3.55, 1.9, 0.7, size=12, color=COLORS["muted"])
        if i < len(s["steps"]) - 1:
            draw_arrow(slide, x + 2.52, 3.35, x + 2.95, 3.35, COLORS["line"], 2)
    add_text(slide, "Training-time work produces one static sparse delta; inference-time work is a masked addition plus optional antidote.", 1.1, 5.65, 10.7, 0.5, size=17, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_flow(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    x = 0.9
    for i, label in enumerate(s["left"]):
        add_panel(slide, x + i * 2.45, 2.2, 2.05, 0.95, fill=COLORS["white"], line=COLORS["blue"])
        add_text(slide, label, x + i * 2.45 + 0.12, 2.52, 1.82, 0.2, size=12, bold=True, align=PP_ALIGN.CENTER)
        if i < len(s["left"]) - 1:
            draw_arrow(slide, x + i * 2.45 + 2.1, 2.68, x + (i + 1) * 2.45 - 0.05, 2.68, COLORS["line"], 2)
    draw_arrow(slide, 6.9, 2.68, 8.1, 1.75, COLORS["red"], 2.3)
    draw_arrow(slide, 6.9, 2.68, 8.1, 3.65, COLORS["green"], 2.3)
    for i, (cond, action, result) in enumerate(s["branches"]):
        y = 1.15 if i == 0 else 3.05
        color = COLORS["red"] if i == 0 else COLORS["green"]
        add_panel(slide, 8.1, y, 3.55, 1.35, fill=COLORS["white"], line=color)
        add_text(slide, cond, 8.35, y + 0.18, 3.05, 0.24, size=14, bold=True, color=color)
        add_text(slide, action, 8.35, y + 0.55, 3.05, 0.24, size=12, color=COLORS["ink"])
        add_text(slide, result, 8.35, y + 0.9, 3.05, 0.24, size=12, color=COLORS["muted"])
    add_footer(slide)


def slide_code_map(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    y = 1.35
    for title, path in s["items"]:
        add_panel(slide, 1.0, y, 11.0, 0.62, fill=COLORS["white"], line=COLORS["line"])
        add_text(slide, title, 1.25, y + 0.18, 3.3, 0.16, size=12, bold=True, color=COLORS["blue"])
        add_text(slide, path, 4.95, y + 0.18, 6.5, 0.16, size=12, color=COLORS["dark"])
        y += 0.82
    add_text(slide, "For a live artifact discussion, start from README.md, then use reproduce.py --exp1 for the main table.", 1.2, 6.1, 10.5, 0.35, size=15, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_setup(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_bullets(slide, s["bullets"], 0.9, 1.35, 6.2, 4.9, size=17, gap=5)
    add_panel(slide, 7.65, 1.55, 4.1, 4.35, fill=COLORS["white"], line=COLORS["green"])
    for i, (big, small) in enumerate([("3", "datasets"), ("4+", "architectures"), ("6+", "baseline families"), ("5", "evaluation axes")]):
        y = 1.95 + i * 0.9
        add_text(slide, big, 8.0, y, 0.8, 0.35, size=24, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER)
        add_text(slide, small, 8.95, y + 0.08, 2.35, 0.2, size=14, color=COLORS["muted"])
    add_footer(slide)


def slide_bars(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_text(slide, "Accuracy (%)", 1.0, 1.25, 2.0, 0.25, size=12, color=COLORS["muted"])
    max_v = 100
    base_x, base_y = 1.15, 5.65
    chart_h = 3.8
    bar_w = 1.25
    gap = 0.85
    for i, (label, value, color) in enumerate(s["bars"]):
        h = chart_h * value / max_v
        x = base_x + i * (bar_w + gap)
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(base_y - h), Inches(bar_w), Inches(h))
        set_fill(rect, color)
        rect.line.color.rgb = color
        add_text(slide, f"{value:.1f}", x, base_y - h - 0.35, bar_w, 0.2, size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, x - 0.25, base_y + 0.15, bar_w + 0.5, 0.25, size=12, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_panel(slide, 8.0, 1.55, 3.8, 3.6, fill=COLORS["white"], line=COLORS["green"])
    add_text(slide, "Reading the result", 8.35, 1.9, 3.1, 0.3, size=18, bold=True, color=COLORS["green"])
    add_bullets(slide, [
        "Unauthorized accuracy should be low.",
        "VIP accuracy should stay high.",
        "GSUAP gives the lockout; credential gives the recovery.",
    ], 8.35, 2.45, 3.05, 1.5, size=14, gap=3)
    add_text(slide, s["caption"], 1.0, 6.45, 10.9, 0.42, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_overhead(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    colors = [COLORS["blue"], COLORS["green"], COLORS["orange"]]
    for i, (name, metric, note1, note2) in enumerate(s["metrics"]):
        x = 0.95 + i * 4.0
        add_panel(slide, x, 1.6, 3.4, 4.4, fill=COLORS["white"], line=colors[i])
        add_text(slide, name, x + 0.25, 2.0, 2.9, 0.35, size=18, bold=True, color=colors[i])
        add_text(slide, metric, x + 0.25, 2.75, 2.9, 0.5, size=24, bold=True, color=COLORS["dark"], align=PP_ALIGN.CENTER)
        add_text(slide, note1, x + 0.25, 3.65, 2.9, 0.28, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        add_text(slide, note2, x + 0.25, 4.2, 2.9, 0.28, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_attack_results(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_panel(slide, 0.85, 1.4, 5.6, 4.8, fill=COLORS["white"], line=COLORS["blue"])
    add_panel(slide, 6.85, 1.4, 5.6, 4.8, fill=COLORS["white"], line=COLORS["green"])
    add_text(slide, s["left_title"], 1.15, 1.75, 4.8, 0.3, size=18, bold=True, color=COLORS["blue"])
    add_bullets(slide, s["left"], 1.15, 2.3, 4.65, 2.3, size=16, gap=4)
    add_text(slide, s["right_title"], 7.15, 1.75, 4.8, 0.3, size=18, bold=True, color=COLORS["green"])
    add_bullets(slide, s["right"], 7.15, 2.3, 4.65, 2.3, size=16, gap=4)
    add_text(slide, "Feature-space corruption makes collected soft labels and inversion signals low-value.", 1.2, 5.6, 10.8, 0.35, size=16, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_mini_table(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    x, y = 1.0, 1.55
    col_w = [4.0, 1.5, 2.2, 2.2]
    row_h = 0.7
    for c, header in enumerate(s["columns"]):
        add_panel(slide, x + sum(col_w[:c]), y, col_w[c], row_h, fill=COLORS["dark"], line=COLORS["dark"])
        add_text(slide, header, x + sum(col_w[:c]) + 0.05, y + 0.22, col_w[c] - 0.1, 0.18, size=10, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    for r, row in enumerate(s["rows"]):
        ry = y + row_h * (r + 1)
        fill = RGBColor(250, 252, 254) if r % 2 == 0 else COLORS["white"]
        for c, val in enumerate(row):
            add_panel(slide, x + sum(col_w[:c]), ry, col_w[c], row_h, fill=fill, line=COLORS["line"])
            add_text(slide, val, x + sum(col_w[:c]) + 0.05, ry + 0.22, col_w[c] - 0.1, 0.18, size=12, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_text(slide, "Takeaway: original private data is not a hard requirement; surrogate domains can still initialize a useful lock.", 1.0, 5.7, 10.7, 0.45, size=17, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_four_cards(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    colors = [COLORS["blue"], COLORS["green"], COLORS["purple"], COLORS["orange"]]
    for i, (title, desc) in enumerate(s["cards"]):
        x = 0.9 + (i % 2) * 6.0
        y = 1.45 + (i // 2) * 2.25
        add_panel(slide, x, y, 5.25, 1.75, fill=COLORS["white"], line=colors[i])
        add_text(slide, title, x + 0.25, y + 0.28, 4.6, 0.28, size=17, bold=True, color=colors[i])
        add_text(slide, desc, x + 0.25, y + 0.82, 4.6, 0.42, size=14, color=COLORS["muted"])
    add_footer(slide)


def slide_takeaway(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    add_text(slide, s["headline"], 1.2, 1.45, 10.8, 1.0, size=30, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER)
    for i, point in enumerate(s["points"]):
        y = 3.0 + i * 0.72
        add_panel(slide, 2.15, y, 8.9, 0.48, fill=COLORS["white"], line=COLORS["line"])
        add_text(slide, point, 2.35, y + 0.13, 8.5, 0.15, size=15, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_sources(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, s["title"], idx)
    y = 1.35
    for topic, refs in s["sources"]:
        add_text(slide, topic, 0.95, y, 2.6, 0.22, size=13, bold=True, color=COLORS["blue"])
        add_text(slide, refs, 3.65, y, 8.5, 0.22, size=12, color=COLORS["muted"])
        y += 0.62
    add_footer(slide)


def add_slide(prs, s, idx, images):
    layout = s.get("layout")
    if idx == 1:
        return slide_title(prs, s, idx)
    if layout == "timeline":
        return slide_timeline(prs, s, idx)
    if layout == "two_column":
        return slide_two_column(prs, s, idx)
    if layout == "attack_loop":
        return slide_attack_loop(prs, s, idx)
    if layout == "three_cards":
        return slide_three_cards(prs, s, idx)
    if layout == "matrix":
        return slide_matrix(prs, s, idx)
    if layout == "statement":
        return slide_statement(prs, s, idx)
    if layout == "roles":
        return slide_roles(prs, s, idx)
    if layout == "image":
        return slide_image(prs, s, idx, images)
    if layout == "credential":
        return slide_credential(prs, s, idx)
    if layout == "pipeline":
        return slide_pipeline(prs, s, idx)
    if layout == "flow":
        return slide_flow(prs, s, idx)
    if layout == "code_map":
        return slide_code_map(prs, s, idx)
    if layout == "setup":
        return slide_setup(prs, s, idx)
    if layout == "bars":
        return slide_bars(prs, s, idx)
    if layout == "overhead":
        return slide_overhead(prs, s, idx)
    if layout == "attack_results":
        return slide_attack_results(prs, s, idx)
    if layout == "mini_table":
        return slide_mini_table(prs, s, idx)
    if layout == "four_cards":
        return slide_four_cards(prs, s, idx)
    if layout == "takeaway":
        return slide_takeaway(prs, s, idx)
    if layout == "sources":
        return slide_sources(prs, s, idx)
    raise ValueError(f"Unknown layout: {layout}")


def build_notes() -> None:
    lines = [
        "# LymphNode 17-minute speaker notes",
        "",
        "Target pacing: 23 main slides plus one backup slide. Average pace is about 40-45 seconds per main slide, with extra time on framework and main result slides.",
        "",
    ]
    for i, s in enumerate(SLIDES, start=1):
        lines.append(f"## {i}. {s['title']} ({s.get('time', '')})")
        lines.append(textwrap.fill(s["notes"], width=100))
        lines.append("")
    lines.extend([
        "## Related-work source links used for the background slides",
        "",
        "- Tramer et al., Stealing Machine Learning Models via Prediction APIs: https://arxiv.org/abs/1609.02943",
        "- Orekondy et al., Knockoff Nets: https://arxiv.org/abs/1812.02766",
        "- Fredrikson et al., model inversion attacks: https://dl.acm.org/doi/10.1145/2810103.2813677",
        "- Lukas et al., SoK on DNN watermarking robustness: https://arxiv.org/abs/2108.04974",
        "- Deep-Lock: https://arxiv.org/abs/2008.05966",
        "- Prediction Poisoning: https://arxiv.org/abs/1906.10908",
        "- AdvParams: https://arxiv.org/abs/2105.13697",
        "- ActiveGuard: https://arxiv.org/abs/2103.01527",
        "- GD-UAP: https://arxiv.org/abs/1801.08092",
    ])
    NOTES_OUT.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    images = render_pdf_figures()
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for idx, s in enumerate(SLIDES, start=1):
        add_slide(prs, s, idx, images)

    prs.save(PPTX_OUT)
    build_notes()
    print(PPTX_OUT)
    print(NOTES_OUT)


if __name__ == "__main__":
    main()
